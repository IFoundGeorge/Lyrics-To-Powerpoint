import os
import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import pptx
from pptx.enum.text import MSO_ANCHOR
from tkinter import filedialog
from pptx.dml.color import RGBColor
from tkinter import messagebox

# Create a splash screen with program details
splash_screen = tk.Toplevel()
splash_screen.title("Lyrics to PowerPoint")
splash_screen.geometry("450x250")
splash_screen.attributes("-alpha", 0.9)

title_label = ttk.Label(splash_screen, text="Lyrics to PowerPoint")
title_label.config(font=("Helvetica", 28))
title_label.pack(pady=(50, 0))

subtitle_label = ttk.Label(splash_screen, text="A program for Praise and Worship Team\nDeveloped by Cyrus Cavero | using ChatGPT and PyCharm")
subtitle_label.config(font=("Helvetica", 12))
subtitle_label.pack(pady=(10, 20))

# Function to close the splash screen and start the main program
def start_program():
    splash_screen.destroy()
    root.deiconify()

# Create the main window and its widgets
root = tk.Tk()
root.withdraw()
root.title("Lyrics to PowerPoint")
root.geometry("400x300")

lyrics_label = ttk.Label(root, text="Enter the lyrics:")
lyrics_label.pack(pady=(10, 5))

lyrics_input = tk.Text(root, height=10)
lyrics_input.pack(padx=10)

info_label = ttk.Label(root, text="Note: If the line of the lyrics is long\nplease cut half of the line so that it won't overlap the slide page.")
info_label.pack(pady=(5, 10))

button_frame = ttk.Frame(root)
button_frame.pack(pady=10)


def create_powerpoint(lyrics):
    # Split the lyrics into sentences
    sentences = lyrics.split("\n")

    # Create a new presentation
    prs = Presentation()

    # Add a new slide with a centered title for each pair of sentences
    slide_layout = prs.slide_layouts[6]
    for i in range(0, len(sentences), 2):
        first_sentence = sentences[i].strip()
        if i + 1 < len(sentences):
            second_sentence = sentences[i + 1].strip()
        else:
            second_sentence = ""

        # Create a new slide
        slide = prs.slides.add_slide(slide_layout)

        # Set the text box to the center of the slide
        textbox_shape = slide.shapes.add_textbox(left=Inches(1.25), top=Inches(1.5), width=Inches(7.5),
                                                 height=Inches(5))

        # Add the first sentence to the text box
        p = textbox_shape.text_frame.add_paragraph()
        p.text = first_sentence
        p.font.size = Pt(60)
        p.font.bold = False
        p.alignment = pptx.enum.text.PP_ALIGN.CENTER
        p.vertical_anchor = MSO_ANCHOR.MIDDLE
        p.word_wrap = True

        # Add the second sentence to the text box
        if second_sentence != "":
            p2 = textbox_shape.text_frame.add_paragraph()
            p2.text = second_sentence
            p2.font.size = Pt(60)
            p2.font.bold = False
            p2.alignment = pptx.enum.text.PP_ALIGN.CENTER
            p2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2.word_wrap = True
            if p.font.size == Pt(60) and not p.font.bold:
                p2.font.size = Pt(60)
                p2.font.bold = False

    # Ask user to select the file name and save location
    file_name = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint file", "*.pptx")])

    # Save the presentation
    try:
        prs.save(file_name)
        messagebox.showinfo("Success!", "The lyrics are successfully built as a Pptx.")
    except:
        messagebox.showerror("Error!", "The lyrics are not successfully built as a Pptx. Please contact the developer.")
        return

create_powerpoint_button = ttk.Button(button_frame, text="Create PowerPoint", command=lambda: create_powerpoint(lyrics_input.get("1.0", tk.END)))
create_powerpoint_button.pack(side="left", padx=10)

# Show the splash screen and start the program after 3 seconds
splash_screen.after(3000, start_program)

root.mainloop()
